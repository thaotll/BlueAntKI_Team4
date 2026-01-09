"""
Renders PptxPresentationModel to actual PPTX bytes using python-pptx.
Pure rendering logic - no business/content decisions.
"""

import logging
from io import BytesIO
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from app.models.pptx import (
    PptxPresentationModel,
    PptxSlideModel,
    TextBoxShape,
    ImageShape,
    RectangleShape,
    ChartShape,
    TableShape,
    ShapeModel,
    TextParagraph,
    TextStyle,
    RgbColor,
    SlideLayout,
)

logger = logging.getLogger(__name__)


class PptxRenderer:
    """
    Renders presentation models to PPTX files.

    Usage:
        renderer = PptxRenderer()
        pptx_bytes = renderer.render(model)
    """

    # Layout index mapping
    LAYOUT_MAP = {
        SlideLayout.BLANK: 6,
        SlideLayout.TITLE: 0,
        SlideLayout.TITLE_CONTENT: 1,
    }

    # Alignment mapping
    ALIGNMENT_MAP = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }

    def render(self, model: PptxPresentationModel) -> bytes:
        """
        Render presentation model to PPTX bytes.

        Args:
            model: The presentation model to render

        Returns:
            PPTX file as bytes
        """
        logger.info(f"Rendering presentation: {model.title}")

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(model.width_inches)
        prs.slide_height = Inches(model.height_inches)

        # Render each slide
        for i, slide_model in enumerate(model.slides):
            logger.debug(f"Rendering slide {i + 1}/{len(model.slides)}")
            self._render_slide(prs, slide_model)

        # Save to bytes
        stream = BytesIO()
        prs.save(stream)
        stream.seek(0)

        pptx_bytes = stream.getvalue()
        logger.info(f"Rendered {len(model.slides)} slides ({len(pptx_bytes)} bytes)")

        return pptx_bytes

    def _render_slide(self, prs: Presentation, slide_model: PptxSlideModel) -> None:
        """Render a single slide."""
        layout_idx = self.LAYOUT_MAP.get(slide_model.layout, 6)
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Render all shapes
        for shape_model in slide_model.shapes:
            self._render_shape(slide, shape_model)

        # Add speaker notes if present
        if slide_model.notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_model.notes

    def _render_shape(self, slide, shape_model: ShapeModel) -> None:
        """Render a shape based on its type."""
        if isinstance(shape_model, TextBoxShape):
            self._render_textbox(slide, shape_model)
        elif isinstance(shape_model, ImageShape):
            self._render_image(slide, shape_model)
        elif isinstance(shape_model, RectangleShape):
            self._render_rectangle(slide, shape_model)
        elif isinstance(shape_model, ChartShape):
            self._render_chart(slide, shape_model)
        elif isinstance(shape_model, TableShape):
            self._render_table(slide, shape_model)

    def _render_textbox(self, slide, model: TextBoxShape) -> None:
        """Render a text box shape."""
        pos = model.position
        textbox = slide.shapes.add_textbox(
            Inches(pos.x),
            Inches(pos.y),
            Inches(pos.width),
            Inches(pos.height),
        )

        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        # Render paragraphs
        for i, para_model in enumerate(model.paragraphs):
            if i == 0:
                para = text_frame.paragraphs[0]
            else:
                para = text_frame.add_paragraph()

            self._render_paragraph(para, para_model, model.default_style)

    def _render_paragraph(
        self,
        para,
        para_model: TextParagraph,
        default_style: Optional[TextStyle],
    ) -> None:
        """Render a paragraph with its text runs."""
        para.alignment = self.ALIGNMENT_MAP.get(para_model.alignment, PP_ALIGN.LEFT)
        para.space_after = Pt(para_model.space_after_pt)

        for i, run_model in enumerate(para_model.runs):
            if i == 0:
                run = para.runs[0] if para.runs else para.add_run()
            else:
                run = para.add_run()

            run.text = run_model.text

            # Apply styling (run style > default style)
            style = run_model.style or default_style
            if style:
                self._apply_text_style(run, style)

    def _apply_text_style(self, run, style: TextStyle) -> None:
        """Apply text styling to a run."""
        run.font.name = style.font_name
        run.font.size = Pt(style.font_size_pt)
        run.font.bold = style.bold
        run.font.italic = style.italic
        if style.color:
            run.font.color.rgb = self._to_rgb_color(style.color)

    def _render_image(self, slide, model: ImageShape) -> None:
        """Render an image shape."""
        pos = model.position
        if model.image_path:
            try:
                slide.shapes.add_picture(
                    model.image_path,
                    Inches(pos.x),
                    Inches(pos.y),
                    width=Inches(pos.width),
                )
            except Exception as e:
                logger.warning(f"Failed to add image {model.image_path}: {e}")
        elif model.image_bytes:
            try:
                from io import BytesIO
                image_stream = BytesIO(model.image_bytes)
                slide.shapes.add_picture(
                    image_stream,
                    Inches(pos.x),
                    Inches(pos.y),
                    width=Inches(pos.width),
                )
            except Exception as e:
                logger.warning(f"Failed to add image from bytes: {e}")

    def _render_rectangle(self, slide, model: RectangleShape) -> None:
        """Render a rectangle shape."""
        pos = model.position
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(pos.x),
            Inches(pos.y),
            Inches(pos.width),
            Inches(pos.height),
        )

        if model.fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._to_rgb_color(model.fill_color)
        else:
            shape.fill.background()

        if model.line_color:
            shape.line.color.rgb = self._to_rgb_color(model.line_color)
            shape.line.width = Pt(model.line_width_pt)
        else:
            shape.line.fill.background()

    def _render_chart(self, slide, model: ChartShape) -> None:
        """
        Render a chart shape.
        
        Charts are pre-rendered as images by the chart_generator,
        so we render them as image shapes using the rendered_image_bytes.
        """
        pos = model.position
        
        # Charts should have been pre-rendered to image bytes
        if model.rendered_image_bytes:
            try:
                image_stream = BytesIO(model.rendered_image_bytes)
                slide.shapes.add_picture(
                    image_stream,
                    Inches(pos.x),
                    Inches(pos.y),
                    width=Inches(pos.width),
                )
            except Exception as e:
                logger.warning(f"Failed to add chart image: {e}")
        else:
            # Fallback: create a placeholder rectangle
            logger.warning("Chart has no rendered image, creating placeholder")
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(pos.x),
                Inches(pos.y),
                Inches(pos.width),
                Inches(pos.height),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
            
            # Add placeholder text
            text_frame = shape.text_frame
            text_frame.text = f"[Chart: {model.chart_type.value}]"

    def _render_table(self, slide, model: TableShape) -> None:
        """Render a table shape."""
        pos = model.position
        
        if not model.rows:
            logger.warning("Table has no rows, skipping")
            return
        
        # Determine table dimensions
        row_count = len(model.rows)
        col_count = max(len(row.cells) for row in model.rows) if model.rows else 1
        
        try:
            # Add table to slide
            table_shape = slide.shapes.add_table(
                row_count,
                col_count,
                Inches(pos.x),
                Inches(pos.y),
                Inches(pos.width),
                Inches(pos.height),
            )
            table = table_shape.table
            
            # Calculate column widths (equal distribution)
            col_width = Inches(pos.width / col_count)
            for i in range(col_count):
                table.columns[i].width = col_width
            
            # Render each row
            for row_idx, row_model in enumerate(model.rows):
                row = table.rows[row_idx]
                
                for col_idx, cell_model in enumerate(row_model.cells):
                    if col_idx >= col_count:
                        break
                    
                    cell = row.cells[col_idx]
                    cell.text = cell_model.text
                    
                    # Apply cell styling
                    para = cell.text_frame.paragraphs[0]
                    
                    # Determine style (cell style > row style > table defaults)
                    if row_model.is_header and model.header_style:
                        style = cell_model.style or model.header_style
                    else:
                        style = cell_model.style or model.cell_style
                    
                    if style:
                        if para.runs:
                            run = para.runs[0]
                        else:
                            run = para.add_run()
                            run.text = cell_model.text
                            cell.text_frame.paragraphs[0].clear()
                        
                        run.font.name = style.font_name
                        run.font.size = Pt(style.font_size_pt)
                        run.font.bold = style.bold
                        if style.color:
                            run.font.color.rgb = self._to_rgb_color(style.color)
                    
                    # Apply cell background color
                    if row_model.is_header and model.header_background:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = self._to_rgb_color(model.header_background)
                    elif cell_model.background_color:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = self._to_rgb_color(cell_model.background_color)
                    elif model.alternating_row_color and row_idx % 2 == 1:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = self._to_rgb_color(model.alternating_row_color)
        
        except Exception as e:
            logger.error(f"Failed to render table: {e}", exc_info=True)

    @staticmethod
    def _to_rgb_color(color: RgbColor) -> RGBColor:
        """Convert model RgbColor to python-pptx RGBColor."""
        return RGBColor(color.r, color.g, color.b)
